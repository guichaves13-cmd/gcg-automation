"""
PPT to Video — Convert PowerPoint presentations to video with narration.
Extracts slide images and text, generates voiceover, and compiles to video.
"""
import os
import subprocess
import tempfile
import json
import re
from pathlib import Path


def extract_slides(pptx_path, output_dir, dpi=150):
    """Extract each slide as an image. Returns list of slide image paths."""
    os.makedirs(output_dir, exist_ok=True)
    slides = []

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return None, "python-pptx not installed. pip install python-pptx"

    try:
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)

        # Try using LibreOffice for slide export
        try:
            subprocess.run([
                "soffice", "--headless", "--convert-to", "pdf",
                "--outdir", output_dir, pptx_path
            ], capture_output=True, text=True, timeout=60)
            pdf_path = os.path.join(output_dir, Path(pptx_path).stem + ".pdf")
            if os.path.exists(pdf_path):
                # Convert PDF pages to images
                from pdf2image import convert_from_path
                images = convert_from_path(pdf_path, dpi=dpi)
                for i, img in enumerate(images):
                    slide_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
                    img.save(slide_path, "PNG")
                    text = ""
                    if i < slide_count:
                        texts = []
                        for shape in prs.slides[i].shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    t = para.text.strip()
                                    if t:
                                        texts.append(t)
                        text = "\n".join(texts)
                    slides.append({
                        "index": i + 1, "path": slide_path,
                        "text": text, "total": slide_count,
                    })
                try: os.unlink(pdf_path)
                except: pass
                return slides
        except FileNotFoundError:
            pass
        except Exception:
            pass

        # Fallback: extract text only, generate placeholder slides
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            slide_text = "\n".join(texts)
            # Generate a colored background slide with FFmpeg
            slide_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
            _generate_text_slide(slide_text[:200] or f"Slide {i+1}",
                                 slide_path, width=1920, height=1080)
            slides.append({
                "index": i + 1, "path": slide_path,
                "text": slide_text, "total": slide_count,
            })

        return slides

    except Exception as e:
        return None, str(e)


def _generate_text_slide(text, output_path, width=1920, height=1080):
    """Generate a text slide image using FFmpeg drawtext."""
    safe_text = text.replace("'", "\\'").replace(":", "\\:").replace("\n", "\\n")
    safe_text = safe_text[:200]
    cmd = [
        "ffmpeg", "-y",
        "-f", "lavfi", "-i", f"color=c=#1a1a2e:s={width}x{height}:d=3",
        "-vf", f"drawtext=text='{safe_text}':fontcolor=white:fontsize=36:x=(w-text_w)/2:y=(h-text_h)/2:fontfile=Arial:line_spacing=10",
        "-frames:v", "1",
        output_path,
    ]
    subprocess.run(cmd, capture_output=True, text=True, timeout=30)


def generate_narration(slide_texts, language="portuguese", output_path=None):
    """
    Generate AI narration script from slide texts and convert to speech.
    Returns path to narration audio.
    """
    from core.ai_engine import ask_ai

    combined = "\n\n".join(
        f"Slide {s['index']}: {s['text'][:300]}"
        for s in slide_texts if s.get("text", "").strip()
    )

    lang = "português brasileiro" if language == "portuguese" else "english"
    prompt = f"""You are a video narrator. Write a natural, conversational narration script for the following slides.

Language: {lang}
Style: Professional, clear, engaging
Rules:
- Write as连贯 narration, not "slide 1, slide 2"
- Max 30 seconds per slide
- Plain paragraphs only

Slides content:
{combined[:10000]}

Narration:
"""
    try:
        script = ask_ai(prompt)
        if script.startswith("["):
            raise Exception(script)

        import edge_tts
        import asyncio

        if not output_path:
            fd, output_path = tempfile.mkstemp(suffix=".mp3")
            os.close(fd)

        voice_map = {"portuguese": "pt-BR-AntonioNeural", "english": "en-US-JennyNeural"}
        voice = voice_map.get(language, "en-US-JennyNeural")

        async def _gen():
            communicate = edge_tts.Communicate(script, voice)
            await communicate.save(output_path)

        asyncio.run(_gen())
        return output_path, script
    except Exception as e:
        return None, str(e)


def ppt_to_video(pptx_path, output_path, language="portuguese",
                 slide_duration=8, transition="fade", ffmpeg_path="ffmpeg"):
    """
    Full PPT to video conversion: extract slides → generate narration → create video.
    """
    tmp_dir = tempfile.mkdtemp(prefix="ppt2vid_")
    try:
        slides = extract_slides(pptx_path, tmp_dir)
        if isinstance(slides, tuple) and slides[0] is None:
            raise RuntimeError(slides[1])

        # Generate narration
        audio_path, script = generate_narration(slides, language)
        if not audio_path or not os.path.exists(audio_path):
            # No narration: use per-slide timing
            concat_file = os.path.join(tmp_dir, "concat.txt")
            with open(concat_file, "w") as f:
                for slide in slides:
                    dur = slide_duration if slide.get("text", "").strip() else 4
                    f.write(f"file '{slide['path']}'\n")
                    f.write(f"duration {dur}\n")

            cmd = [
                ffmpeg_path, "-y", "-f", "concat", "-safe", "0",
                "-i", concat_file,
                "-c:v", "libx264", "-preset", "medium", "-crf", "20",
                "-pix_fmt", "yuv420p",
                "-r", "30",
                output_path,
            ]
            subprocess.run(cmd, capture_output=True, text=True, timeout=300, check=True)
        else:
            # Use narration audio to time slides — derive ffprobe path safely
            _ff_dir = os.path.dirname(ffmpeg_path)
            _ff_base = os.path.basename(ffmpeg_path)
            _probe_base = _ff_base.replace("ffmpeg", "ffprobe")
            _ffprobe = os.path.join(_ff_dir, _probe_base) if _ff_dir else _probe_base
            if not os.path.isfile(_ffprobe):
                import shutil as _sh
                _ffprobe = _sh.which("ffprobe") or "ffprobe"
            result = subprocess.run(
                [_ffprobe, "-v", "quiet",
                 "-show_entries", "format=duration", "-of", "csv=p=0", audio_path],
                capture_output=True, text=True, timeout=15
            )
            audio_dur = float(result.stdout.strip() or 0)
            per_slide_dur = max(3, audio_dur / max(len(slides), 1))

            concat_file = os.path.join(tmp_dir, "concat.txt")
            with open(concat_file, "w") as f:
                for slide in slides:
                    f.write(f"file '{slide['path']}'\n")
                    f.write(f"duration {per_slide_dur}\n")

            cmd = [
                ffmpeg_path, "-y", "-f", "concat", "-safe", "0",
                "-i", concat_file,
                "-i", audio_path,
                "-c:v", "libx264", "-preset", "medium", "-crf", "20",
                "-c:a", "aac", "-b:a", "192k",
                "-pix_fmt", "yuv420p",
                "-shortest",
                "-r", "30",
                output_path,
            ]
            subprocess.run(cmd, capture_output=True, text=True, timeout=300, check=True)

        return output_path
    finally:
        import shutil
        shutil.rmtree(tmp_dir, ignore_errors=True)
